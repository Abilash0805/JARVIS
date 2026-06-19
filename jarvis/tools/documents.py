"""Document & artifact builders: PowerPoint, PDF, and multi-page websites.

These let JARVIS produce finished deliverables in one shot instead of dictating
raw bytes through ``write_file`` (which can't emit binary .pptx/.pdf). Heavy
deps are imported lazily so the package stays import-safe without them; install
with ``pip install -e ".[documents]"`` (python-pptx + reportlab).

Inputs are plain JSON-friendly structures so any model can drive them:

- ``create_pptx``: slides = [{"title": str, "bullets": [str, ...], "notes"?: str}]
- ``create_pdf``:  blocks = [{"type": "heading"|"subheading"|"text"|"bullets",
                              "text"|"items": ...}]
- ``create_website``: pages = [{"filename": str, "title": str, "body_html": str}]
"""

from __future__ import annotations

import json
import os
from pathlib import Path
from typing import Any

from jarvis.tools.base import Tool, ToolError
from jarvis.utils.safety import SafetyGate


def _resolve(path: str) -> Path:
    return Path(os.path.expanduser(path)).resolve()


def _coerce(value: Any, what: str) -> Any:
    """Accept either a real list/dict or a JSON string of one (models vary)."""
    if isinstance(value, str):
        try:
            return json.loads(value)
        except json.JSONDecodeError as exc:
            raise ToolError(f"{what} must be JSON; could not parse: {exc}") from exc
    return value


# --------------------------------------------------------------------------- #
# PowerPoint — hand-drawn theme (not the bland default template)
# --------------------------------------------------------------------------- #
# (background, heading/accent-bar, title text, body text, muted/footer text)
_PPTX_THEMES = {
    "professional": ("FFFFFF", "1F3A5F", "1F3A5F", "1A1A2E", "6B7280"),
    "modern":       ("FFFFFF", "0F766E", "0F766E", "1F2937", "6B7280"),
    "dark":         ("111827", "1F2937", "38BDF8", "E5E7EB", "9CA3AF"),
    "minimal":      ("FFFFFF", "111111", "111111", "111111", "6B7280"),
}
_PPTX_ACCENTS = {  # thin highlight bar / bullet glyph color per theme
    "professional": "C9A227",
    "modern": "FF6B6B",
    "dark": "38BDF8",
    "minimal": "2563EB",
}


def _build_pptx(
    path: str, title: str, slides: Any, subtitle: str = "", theme: str = "professional"
) -> str:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.util import Emu, Inches, Pt
    except ImportError as exc:  # pragma: no cover - dep missing
        raise ToolError(
            "python-pptx not installed. Run: pip install python-pptx"
        ) from exc

    slides = _coerce(slides, "slides")
    if not isinstance(slides, list):
        raise ToolError("slides must be a list of {title, bullets} objects")

    bg_hex, bar_hex, title_hex, body_hex, muted_hex = _PPTX_THEMES.get(
        theme, _PPTX_THEMES["professional"]
    )
    accent_hex = _PPTX_ACCENTS.get(theme, _PPTX_ACCENTS["professional"])
    rgb = lambda h: RGBColor.from_string(h)  # noqa: E731

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W, H = prs.slide_width, prs.slide_height

    def fill_bg(slide, color_hex: str) -> None:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(color_hex)

    def add_rect(slide, left, top, width, height, color_hex: str):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(color_hex)
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def add_text(
        slide, left, top, width, height, text, *, size, color_hex, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    ):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = rgb(color_hex)
        return box

    total = len(slides) + 1

    # --- Title slide: full-bleed color panel with an accent rule. ---
    s = prs.slides.add_slide(blank)
    fill_bg(s, bar_hex)
    add_rect(s, 0, H - Emu(int(H * 0.012)), W, Emu(int(H * 0.012)), accent_hex)
    add_text(
        s, Inches(1), Inches(2.7), W - Inches(2), Inches(1.4), title,
        size=44, color_hex="FFFFFF", bold=True, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    if subtitle:
        add_text(
            s, Inches(1), Inches(4.0), W - Inches(2), Inches(0.8), subtitle,
            size=20, color_hex="E5E7EB", align=PP_ALIGN.CENTER,
        )

    # --- Content slides: accent bar, themed title, bulleted body, footer. ---
    for i, slide in enumerate(slides, 1):
        if not isinstance(slide, dict):
            raise ToolError(f"slide {i} must be an object with title/bullets")
        s = prs.slides.add_slide(blank)
        fill_bg(s, bg_hex)
        add_rect(s, 0, 0, W, Inches(0.12), bar_hex)
        add_text(
            s, Inches(0.7), Inches(0.4), W - Inches(1.4), Inches(0.9),
            str(slide.get("title", f"Slide {i}")),
            size=30, color_hex=title_hex, bold=True,
        )
        add_rect(s, Inches(0.7), Inches(1.25), Inches(1.0), Pt(3), accent_hex)

        bullets = slide.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        body_box = s.shapes.add_textbox(
            Inches(0.7), Inches(1.6), W - Inches(1.4), H - Inches(2.2)
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        for j, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"•  {bullet}"
            p.font.size = Pt(18)
            p.font.color.rgb = rgb(body_hex)
            p.space_after = Pt(10)

        notes = slide.get("notes")
        if notes:
            s.notes_slide.notes_text_frame.text = str(notes)

        add_text(
            s, Inches(0.7), H - Inches(0.5), Inches(4), Inches(0.35), title,
            size=10, color_hex=muted_hex,
        )
        add_text(
            s, W - Inches(2.2), H - Inches(0.5), Inches(1.5), Inches(0.35),
            f"{i + 1} / {total}", size=10, color_hex=muted_hex, align=PP_ALIGN.RIGHT,
        )

    p = _resolve(path)
    p.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(p))
    return f"created PowerPoint ({theme} theme) with {total} slides at {p}"


# --------------------------------------------------------------------------- #
# PDF — colored header bar, accent headings, page numbers in the footer.
# --------------------------------------------------------------------------- #
_PDF_THEMES = {  # (primary/heading color, accent/bar color, muted footer color)
    "professional": ("#1F3A5F", "#C9A227", "#6B7280"),
    "modern":        ("#0F766E", "#FF6B6B", "#6B7280"),
    "dark":          ("#1F2937", "#38BDF8", "#6B7280"),
    "minimal":       ("#111111", "#2563EB", "#6B7280"),
}


def _build_pdf(path: str, title: str, blocks: Any, theme: str = "professional") -> str:
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.platypus import (
            Flowable,
            ListFlowable,
            ListItem,
            Paragraph,
            SimpleDocTemplate,
            Spacer,
        )
    except ImportError as exc:  # pragma: no cover - dep missing
        raise ToolError("reportlab not installed. Run: pip install reportlab") from exc

    blocks = _coerce(blocks, "blocks")
    if not isinstance(blocks, list):
        raise ToolError("blocks must be a list of {type, text|items} objects")

    primary_hex, accent_hex, muted_hex = _PDF_THEMES.get(
        theme, _PDF_THEMES["professional"]
    )
    primary, accent, muted = (
        colors.HexColor(primary_hex), colors.HexColor(accent_hex), colors.HexColor(muted_hex)
    )

    p = _resolve(path)
    p.parent.mkdir(parents=True, exist_ok=True)
    doc = SimpleDocTemplate(
        str(p), pagesize=letter, title=title,
        topMargin=70, bottomMargin=50, leftMargin=56, rightMargin=56,
    )

    def header_footer(canvas, _doc) -> None:
        canvas.saveState()
        width, height = letter
        canvas.setFillColor(primary)
        canvas.rect(0, height - 14, width, 14, fill=1, stroke=0)
        canvas.setFillColor(accent)
        canvas.rect(0, height - 16, width, 2, fill=1, stroke=0)
        canvas.setFillColor(muted)
        canvas.setFont("Helvetica", 8)
        canvas.drawString(56, 28, title)
        canvas.drawRightString(width - 56, 28, f"Page {canvas.getPageNumber()}")
        canvas.restoreState()

    styles = getSampleStyleSheet()
    heading1 = ParagraphStyle(
        "ThemedHeading1", parent=styles["Heading1"], textColor=primary,
    )
    heading2 = ParagraphStyle(
        "ThemedHeading2", parent=styles["Heading2"], textColor=primary,
    )
    title_style = ParagraphStyle(
        "ThemedTitle", parent=styles["Title"], textColor=primary,
    )
    story: list[Any] = []

    class _AccentRule(Flowable):
        """A thin colored horizontal rule, used under the title."""

        def __init__(self, color, height: float = 2.2, width_fraction: float = 0.18) -> None:
            super().__init__()
            self._color = color
            self.height = height
            self._width_fraction = width_fraction

        def wrap(self, avail_width, avail_height):  # noqa: ARG002 - reportlab API
            self.width = avail_width
            return avail_width, self.height

        def draw(self) -> None:  # noqa: D401 - reportlab API
            self.canv.saveState()
            self.canv.setFillColor(self._color)
            self.canv.rect(
                0, 0, self.width * self._width_fraction, self.height, fill=1, stroke=0
            )
            self.canv.restoreState()

    if title:
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 4))
        story.append(_AccentRule(accent))
        story.append(Spacer(1, 14))

    for blk in blocks:
        if not isinstance(blk, dict):
            raise ToolError("each block must be an object with a 'type'")
        kind = str(blk.get("type", "text")).lower()
        if kind == "heading":
            story.append(Spacer(1, 8))
            story.append(Paragraph(str(blk.get("text", "")), heading1))
        elif kind == "subheading":
            story.append(Paragraph(str(blk.get("text", "")), heading2))
        elif kind in ("bullets", "list"):
            items = blk.get("items") or []
            if isinstance(items, str):
                items = [items]
            story.append(
                ListFlowable(
                    [ListItem(Paragraph(str(it), styles["BodyText"])) for it in items],
                    bulletType="bullet",
                )
            )
        else:  # text / paragraph
            story.append(Paragraph(str(blk.get("text", "")), styles["BodyText"]))
        story.append(Spacer(1, 6))

    doc.build(story, onFirstPage=header_footer, onLaterPages=header_footer)
    return f"created PDF ({theme} theme) with {len(blocks)} blocks at {p}"


# --------------------------------------------------------------------------- #
# Website — themed CSS (futuristic by default, plus clean/dark/minimal)
# --------------------------------------------------------------------------- #
_WEB_FUTURISTIC_CSS = """\
:root{--bg:#05070d;--panel:rgba(13,20,33,.72);--line:rgba(56,189,248,.18);
  --cyan:#38e8ff;--teal:#5eead4;--ice:#bce7ff;--fg:#d7e6f7;--muted:#6f86a8;
  --glow:0 0 22px rgba(56,232,255,.40);}
*{box-sizing:border-box;}
body{margin:0;color:var(--fg);background:var(--bg);line-height:1.65;
  font-family:"Segoe UI",-apple-system,Roboto,Helvetica,Arial,sans-serif;}
body::before{content:"";position:fixed;inset:0;z-index:-1;
  background:radial-gradient(1100px 640px at 80% -10%,rgba(56,232,255,.12),transparent 60%),
    radial-gradient(860px 560px at -6% 110%,rgba(94,234,212,.10),transparent 55%),var(--bg);}
body::after{content:"";position:fixed;inset:0;z-index:-1;opacity:.5;
  background-image:linear-gradient(rgba(56,189,248,.05) 1px,transparent 1px),
    linear-gradient(90deg,rgba(56,189,248,.05) 1px,transparent 1px);
  background-size:46px 46px;
  -webkit-mask-image:radial-gradient(circle at 50% 30%,#000 30%,transparent 85%);
  mask-image:radial-gradient(circle at 50% 30%,#000 30%,transparent 85%);}
header{position:sticky;top:0;z-index:10;display:flex;align-items:center;gap:1.25rem;
  padding:1rem 2rem;background:rgba(8,14,24,.72);backdrop-filter:blur(10px);
  border-bottom:1px solid var(--line);}
header nav a{color:var(--ice);margin-right:1.25rem;text-decoration:none;font-size:.95rem;
  letter-spacing:.5px;opacity:.85;transition:color .2s,text-shadow .2s;}
header nav a:hover{opacity:1;color:var(--cyan);text-shadow:0 0 12px rgba(56,232,255,.6);}
main{max-width:980px;margin:0 auto;padding:3rem 2rem;}
h1,h2,h3{line-height:1.2;color:#eaf9ff;letter-spacing:.5px;
  text-shadow:0 0 18px rgba(56,232,255,.25);}
h1{font-size:2.4rem;}
a{color:var(--cyan);}
p{color:#c2d3e8;}
.card{border:1px solid var(--line);border-radius:16px;padding:1.5rem;margin:1.25rem 0;
  background:var(--panel);backdrop-filter:blur(8px);
  box-shadow:0 8px 28px rgba(0,0,0,.35);transition:transform .18s,border-color .2s,box-shadow .2s;}
.card:hover{transform:translateY(-3px);border-color:rgba(56,232,255,.45);
  box-shadow:0 12px 36px rgba(0,0,0,.45),var(--glow);}
footer{color:var(--muted);text-align:center;padding:2.5rem;font-size:.9rem;
  border-top:1px solid var(--line);margin-top:2rem;}
button,.btn{background:linear-gradient(135deg,#7ef0ff,#38e8ff 45%,#22b8d6);color:#021018;
  border:0;padding:.7rem 1.4rem;border-radius:12px;cursor:pointer;font-weight:700;
  letter-spacing:1px;text-transform:uppercase;text-decoration:none;display:inline-block;
  box-shadow:var(--glow);transition:transform .12s,filter .2s;}
button:hover,.btn:hover{transform:translateY(-1px);filter:brightness(1.08);}
"""

_WEB_CLEAN_CSS = """\
:root{--fg:#1a1a2e;--muted:#555;--accent:#0066ff;--bg:#ffffff;--soft:#f6f8fc;}
*{box-sizing:border-box;}
body{margin:0;font-family:-apple-system,Segoe UI,Roboto,Helvetica,Arial,sans-serif;
  color:var(--fg);background:var(--bg);line-height:1.65;}
header{position:sticky;top:0;background:#fff;border-bottom:1px solid #e7ebf2;
  padding:1rem 2rem;box-shadow:0 1px 8px rgba(0,0,0,.04);}
header nav a{color:var(--fg);margin-right:1.4rem;text-decoration:none;font-weight:600;opacity:.85;}
header nav a:hover{opacity:1;color:var(--accent);}
main{max-width:960px;margin:0 auto;padding:2.5rem 2rem;}
h1,h2,h3{line-height:1.25;}
a{color:var(--accent);}
.card{border:1px solid #e7ebf2;border-radius:14px;padding:1.4rem;margin:1.1rem 0;
  background:var(--soft);box-shadow:0 1px 3px rgba(0,0,0,.05);transition:transform .15s,box-shadow .2s;}
.card:hover{transform:translateY(-2px);box-shadow:0 8px 24px rgba(0,0,0,.08);}
footer{color:var(--muted);text-align:center;padding:2rem;font-size:.9rem;border-top:1px solid #eef1f6;}
button,.btn{background:var(--accent);color:#fff;border:0;padding:.65rem 1.2rem;border-radius:10px;
  cursor:pointer;text-decoration:none;display:inline-block;font-weight:600;}
button:hover,.btn:hover{filter:brightness(1.05);}
"""

_WEB_DARK_CSS = """\
:root{--fg:#e5e7eb;--muted:#9ca3af;--accent:#38bdf8;--bg:#0b0f17;--panel:#111826;}
*{box-sizing:border-box;}
body{margin:0;font-family:-apple-system,Segoe UI,Roboto,Helvetica,Arial,sans-serif;
  color:var(--fg);background:var(--bg);line-height:1.65;}
header{position:sticky;top:0;background:var(--panel);border-bottom:1px solid #1e293b;padding:1rem 2rem;}
header nav a{color:var(--fg);margin-right:1.4rem;text-decoration:none;opacity:.85;}
header nav a:hover{opacity:1;color:var(--accent);}
main{max-width:960px;margin:0 auto;padding:2.5rem 2rem;}
h1,h2,h3{line-height:1.25;color:#f1f5f9;}
a{color:var(--accent);}
.card{border:1px solid #1e293b;border-radius:14px;padding:1.4rem;margin:1.1rem 0;background:var(--panel);}
footer{color:var(--muted);text-align:center;padding:2rem;font-size:.9rem;border-top:1px solid #1e293b;}
button,.btn{background:var(--accent);color:#041018;border:0;padding:.65rem 1.2rem;border-radius:10px;
  cursor:pointer;text-decoration:none;display:inline-block;font-weight:700;}
"""

_WEB_MINIMAL_CSS = """\
:root{--fg:#111;--muted:#666;--accent:#111;--bg:#fff;}
*{box-sizing:border-box;}
body{margin:0;font-family:Georgia,"Times New Roman",serif;color:var(--fg);background:var(--bg);line-height:1.7;}
header{padding:1.5rem 2rem;border-bottom:1px solid #eee;}
header nav a{color:var(--fg);margin-right:1.4rem;text-decoration:none;border-bottom:1px solid transparent;}
header nav a:hover{border-bottom-color:var(--fg);}
main{max-width:720px;margin:0 auto;padding:3rem 2rem;}
h1,h2,h3{line-height:1.25;font-weight:600;}
a{color:var(--accent);}
.card{border-top:1px solid #eee;padding:1.25rem 0;margin:1rem 0;}
footer{color:var(--muted);text-align:center;padding:2.5rem;font-size:.85rem;}
button,.btn{background:var(--fg);color:#fff;border:0;padding:.6rem 1.2rem;border-radius:2px;
  cursor:pointer;text-decoration:none;display:inline-block;}
"""

_WEB_THEMES = {
    "futuristic": _WEB_FUTURISTIC_CSS,
    "clean": _WEB_CLEAN_CSS,
    "dark": _WEB_DARK_CSS,
    "minimal": _WEB_MINIMAL_CSS,
}
_DEFAULT_CSS = _WEB_FUTURISTIC_CSS  # default look matches the JARVIS dashboard

_PAGE_TEMPLATE = """\
<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{title}</title>
<link rel="stylesheet" href="styles.css">
</head>
<body>
<header><nav>{nav}</nav></header>
<main>
{body}
</main>
<footer>{footer}</footer>
</body>
</html>
"""


def _build_website(
    directory: str, pages: Any, site_name: str = "My Site", css: str = "",
    theme: str = "futuristic",
) -> str:
    pages = _coerce(pages, "pages")
    if not isinstance(pages, list) or not pages:
        raise ToolError("pages must be a non-empty list of {filename, title, body_html}")

    root = _resolve(directory)
    root.mkdir(parents=True, exist_ok=True)

    # Build a shared nav from page titles.
    nav_links = []
    for pg in pages:
        fn = str(pg.get("filename", "")).strip()
        if not fn:
            raise ToolError("each page needs a 'filename' (e.g. index.html)")
        nav_links.append(f'<a href="{fn}">{pg.get("title", fn)}</a>')
    nav = "".join(nav_links)
    footer = f"&copy; {site_name}"

    # An explicit `css` override always wins; otherwise pick the theme.
    theme_css = _WEB_THEMES.get(theme, _DEFAULT_CSS)
    (root / "styles.css").write_text(css or theme_css, encoding="utf-8")

    written = ["styles.css"]
    for pg in pages:
        fn = str(pg["filename"]).strip()
        html = _PAGE_TEMPLATE.format(
            title=pg.get("title", site_name),
            nav=nav,
            body=pg.get("body_html", ""),
            footer=footer,
        )
        (root / fn).write_text(html, encoding="utf-8")
        written.append(fn)

    used = "custom" if css else theme
    return (
        f"created website '{site_name}' ({used} theme) in {root} "
        f"({len(written)} files: {', '.join(written)})"
    )


def make_document_tools(gate: SafetyGate) -> list[Tool]:
    """Tools that produce finished artifacts (decks, PDFs, websites)."""

    def create_pptx(
        path: str, title: str, slides: Any, subtitle: str = "",
        theme: str = "professional",
    ) -> str:
        if not gate.confirm(f"CREATE PowerPoint at {path}"):
            raise ToolError("create_pptx denied by safety gate")
        return _build_pptx(path, title, slides, subtitle, theme)

    def create_pdf(path: str, title: str, blocks: Any, theme: str = "professional") -> str:
        if not gate.confirm(f"CREATE PDF at {path}"):
            raise ToolError("create_pdf denied by safety gate")
        return _build_pdf(path, title, blocks, theme)

    def create_website(
        directory: str, pages: Any, site_name: str = "My Site", css: str = "",
        theme: str = "futuristic",
    ) -> str:
        if not gate.confirm(f"CREATE website in {directory}"):
            raise ToolError("create_website denied by safety gate")
        return _build_website(directory, pages, site_name, css, theme)

    _str = {"type": "string"}
    _theme_prop = {
        "type": "string",
        "enum": ["professional", "modern", "dark", "minimal"],
        "description": (
            "visual theme — colored backgrounds/accent bars/footers instead of a "
            "bland default template (default: professional)"
        ),
    }
    return [
        Tool(
            "create_pptx",
            "Build a polished, themed PowerPoint (.pptx) deck — not the bland "
            "default template. 'slides' is a list of objects like "
            '{"title": str, "bullets": [str, ...], "notes": str (optional)}. '
            "A title slide, colored accent bars, and footer page numbers are added "
            "automatically. Choose 'theme': professional, modern, dark, or minimal. "
            "Great for presentations.",
            {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "output .pptx path"},
                    "title": {"type": "string", "description": "deck title"},
                    "subtitle": _str,
                    "theme": _theme_prop,
                    "slides": {
                        "type": "array",
                        "description": "list of {title, bullets[, notes]} slides",
                        "items": {"type": "object"},
                    },
                },
                "required": ["path", "title", "slides"],
            },
            create_pptx,
            dangerous=True,
        ),
        Tool(
            "create_pdf",
            "Build a polished, themed PDF document — colored header bar, accent "
            "rule under the title, and page-numbered footer, not a bare default "
            "layout. 'blocks' is an ordered list of content blocks: "
            '{"type": "heading"|"subheading"|"text"|"bullets", "text": str} or '
            '{"type": "bullets", "items": [str, ...]}. '
            "Choose 'theme': professional, modern, dark, or minimal. "
            "Ideal for reports, study notes, and handouts.",
            {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "output .pdf path"},
                    "title": {"type": "string", "description": "document title"},
                    "theme": _theme_prop,
                    "blocks": {
                        "type": "array",
                        "description": "ordered content blocks",
                        "items": {"type": "object"},
                    },
                },
                "required": ["path", "title", "blocks"],
            },
            create_pdf,
            dangerous=True,
        ),
        Tool(
            "create_website",
            "Scaffold a complete multi-page static website into a directory. "
            "'pages' is a list of {filename, title, body_html}; a shared nav, "
            "responsive styles.css, and footer are generated. Choose 'theme': "
            "futuristic (default — dark glassmorphism with neon-cyan accents, an "
            "animated grid backdrop and glowing cards, matching the JARVIS look), "
            "clean (light & professional), dark, or minimal (serif, editorial). "
            "Pass your own 'css' to fully override the theme.",
            {
                "type": "object",
                "properties": {
                    "directory": {"type": "string", "description": "output folder"},
                    "site_name": _str,
                    "theme": {
                        "type": "string",
                        "enum": ["futuristic", "clean", "dark", "minimal"],
                        "description": "visual theme (default: futuristic)",
                    },
                    "css": {"type": "string", "description": "optional full CSS override"},
                    "pages": {
                        "type": "array",
                        "description": "list of {filename, title, body_html}",
                        "items": {"type": "object"},
                    },
                },
                "required": ["directory", "pages"],
            },
            create_website,
            dangerous=True,
        ),
    ]
